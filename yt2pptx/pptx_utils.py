from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Length, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from tqdm import tqdm

from .video_utils import timestamp_to_seconds


def create_pptx_from_images_with_timestamps(
    image_paths: list[Path],
    output_pptx: Path,
    video_id: str,
    *,
    video_path: Path,
) -> None:
    """Create a PowerPoint presentation from a list of images with timestamps.

    This function takes a list of tuples containing image paths and their corresponding
    timestamps in seconds, and creates a PowerPoint presentation where each slide contains
    an image and a hyperlink to the YouTube video at the specified timestamp.

    Args:
        image_paths (list[Path]): A list of paths to the image file.
        output_pptx (Path): The path where the PowerPoint file will be saved.
        video_id (str): The YouTube video ID to create hyperlinks for the timestamps.
    """
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    for img_path in tqdm(image_paths, desc="Creating slides"):
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(
            str(img_path), Inches(0), Inches(0), width=prs.slide_width
        )

        timestamp = img_path.stem.split("_")[-1].replace("-", ":")
        zero_hours = "0:"
        if timestamp.startswith(zero_hours):
            timestamp = timestamp[len(zero_hours) :]
        seconds = timestamp_to_seconds(timestamp)

        youtube_link = f"https://www.youtube.com/watch?v={video_id}&t={seconds}s"
        youtube_cmd = None
        if video_path:
            bat_path = img_path.with_suffix(".bat")
            bat_path.write_text(
                f"@echo off\n"
                f'ffplay -ss {seconds} -i "{video_path.absolute()}" -x 1920 -loglevel quiet\n'
            )
            youtube_cmd = str(bat_path.absolute())

        left = Inches(0.3)
        top = (
            Length(prs.slide_height - Inches(0.7)) if prs.slide_height else Inches(0.5)
        )
        width = Inches(3)
        height = Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]

        run = p.add_run()
        run.text = f"Jump to {timestamp}"
        font = run.font
        font.size = Pt(12)
        font.bold = True
        font.underline = True
        font.color.rgb = RGBColor(0, 102, 204)
        run.hyperlink.address = youtube_link

        if youtube_cmd:
            btn_height = Inches(0.5)
            btn_width = Inches(0.5)
            btn_top = Length(
                top - btn_height - Inches(0.1)
            )  # Slight spacing above text
            btn = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left,
                btn_top,
                btn_width,
                btn_height,
            )
            btn.text = f"▶"

            # Style text
            text_frame = btn.text_frame
            # Vertical centering
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            p = text_frame.paragraphs[0]
            # Horizontal centering
            p.alignment = PP_ALIGN.CENTER

            p = btn.text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)  # White text
            btn.fill.solid()
            btn.fill.fore_color.rgb = RGBColor(70, 130, 180)  # Steel blue
            btn.line.color.rgb = RGBColor(0, 0, 0)  # Black border

            btn.click_action.hyperlink.address = youtube_cmd

    prs.save(str(output_pptx))
    print(f"✅ PowerPoint saved: {output_pptx}")
