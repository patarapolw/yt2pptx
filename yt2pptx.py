import sys
import subprocess
from pathlib import Path
from PIL import Image
import imagehash
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import yt_dlp
from tqdm import tqdm
import re
import statistics


def sanitize_filename(name):
    return re.sub(r'[\\/*?:"<>|]', "_", name)


def make_timestamp(idx, interval_seconds):
    timestamp = idx * interval_seconds
    hours = timestamp // 3600
    minutes = (timestamp % 3600) // 60
    seconds = timestamp % 60
    return f"{hours:d}-{minutes:02d}-{seconds:02d}"


def sort_timestamp(k: Path):
    timestamp = k.name.split("_").pop()
    if not timestamp:
        return k
    h, x = timestamp.split("-", 1)
    return f"{int(h):04d}-{x}"


def extract_video_id(input_url_or_id: str) -> str:
    # If input is already an 11-char YouTube ID, return it
    if re.fullmatch(r"[A-Za-z0-9_-]{11}", input_url_or_id):
        return input_url_or_id

    # Try to extract video ID from common YouTube URL formats
    patterns = [
        r"(?:v=|\/)([A-Za-z0-9_-]{11})(?:[&?\/]|$)",  # v=ID or /ID
    ]
    for pat in patterns:
        m = re.search(pat, input_url_or_id)
        if m:
            return m.group(1)
    return ""


def download_youtube_video(input_url_or_id: str, out_dir: Path):
    video_id = extract_video_id(input_url_or_id)

    video_folder = out_dir / video_id
    video_folder.mkdir(exist_ok=True)
    title_file = video_folder / ".title.txt"

    final_path = out_dir / f"{video_id}.mp4"

    if final_path.exists() and title_file.exists():
        title = title_file.read_text("utf-8").strip()
        print(f"✅ Video already downloaded: {final_path} with title '{title}'")
        return final_path, title, video_id

    video_info = {}

    def get_info_hook(d):
        if d.get("status") == "finished":
            video_info["title"] = d["info_dict"].get("title", "video")
            video_info["id"] = d["info_dict"].get("id", "")

    ydl_opts = {
        "outtmpl": final_path,
        "format": "bestvideo+bestaudio/best",
        "merge_output_format": "mp4",
        "progress_hooks": [get_info_hook],
        "quiet": True,
    }

    print("🔽 Downloading video...")
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        ydl.download([input_url_or_id])
        info = ydl.extract_info(input_url_or_id, download=True)
        if info is not None:
            video_info["title"] = info.get("title", "video")
            video_info["id"] = info.get("id", "")
        else:
            video_info["title"] = "video"
            video_info["id"] = ""

    title = sanitize_filename(video_info["title"])
    title_file.write_text(title, encoding="utf-8")

    return final_path, title, video_id


def extract_frames_ffmpeg(video_path: Path, frame_dir: Path, interval_seconds=3):
    temp_pattern = frame_dir / "frame_%04d.jpg"
    ffmpeg_cmd = [
        "ffmpeg",
        "-i",
        video_path,
        "-vf",
        f"fps=1/{interval_seconds}",
        "-q:v",
        "2",
        str(temp_pattern),
    ]
    subprocess.run(ffmpeg_cmd, check=True)

    # Rename frames to include timestamp
    frames = sorted(frame_dir.glob("frame_[0-9][0-9][0-9][0-9].jpg"))
    renamed_frames = []
    for idx, frame_path in enumerate(frames):
        timestamp = make_timestamp(idx, interval_seconds)
        new_name = f"frame_{timestamp}.jpg"
        new_path = frame_path.with_name(new_name)
        if new_path.exists():
            new_path.unlink()  # Overwrite if exists

        frame_path.rename(new_path)
        renamed_frames.append(new_path)
    return renamed_frames


def imagehash_to_int(hash_obj: imagehash.ImageHash) -> int:
    bit_string = "".join("1" if b else "0" for b in hash_obj.hash.flatten())
    return int(bit_string, 2)


def filter_unique_images(
    image_paths, hash_diff_threshold: int | None = None, fps_interval=3
):
    hashes = []
    for idx, path in enumerate(tqdm(image_paths, desc="Hashing images")):
        with Image.open(path) as img:
            hash_obj = imagehash.average_hash(img)
            hashes.append((path, hash_obj, idx))

    # Dynamically determine threshold if not provided
    if hash_diff_threshold is None and len(hashes) > 1:
        diffs = [abs(hashes[i][1] - hashes[i - 1][1]) for i in range(1, len(hashes))]
        mean = statistics.mean(diffs)
        stdev = statistics.stdev(diffs) if len(diffs) > 1 else 0
        # Optimal threshold: mean /2, but at least 1
        hash_diff_threshold = max(1, int(mean / 2))
        print(
            f"ℹ️ Auto-calculated hash_diff_threshold: {hash_diff_threshold} using mean/2 with mean={mean:.2f} and stdev={stdev:.2f}"
        )
    elif hash_diff_threshold is None:
        hash_diff_threshold = 5  # fallback for single image
        print(f"ℹ️ Using default hash_diff_threshold: {hash_diff_threshold}")

    unique_images = []
    last_hash = None
    duplicate_start = None
    duplicate_end = None

    for i, (path, curr_hash, idx) in enumerate(hashes):
        if last_hash is None or abs(curr_hash - last_hash) > hash_diff_threshold:
            # Print duplicate range if any
            if duplicate_start is not None and duplicate_end is not None:
                start_sec = (duplicate_start - 1) * fps_interval
                end_sec = duplicate_end * fps_interval
                start_ts = f"{start_sec // 60:02.0f}:{start_sec % 60:02.0f}"
                end_ts = f"{end_sec // 60:02.0f}:{end_sec % 60:02.0f}"
                print(f"🗑️ Removed duplicate frames from {start_ts} to {end_ts}")
                duplicate_start = None
                duplicate_end = None

            seconds = idx * fps_interval
            unique_images.append((path, seconds))
            last_hash = curr_hash
        else:
            if duplicate_start is None:
                duplicate_start = idx
            duplicate_end = idx

    # Print any remaining duplicate range at the end
    if duplicate_start is not None and duplicate_end is not None:
        start_sec = duplicate_start * fps_interval
        end_sec = duplicate_end * fps_interval
        start_ts = f"{start_sec // 60:02.0f}:{start_sec % 60:02.0f}"
        end_ts = f"{end_sec // 60:02.0f}:{end_sec % 60:02.0f}"
        print(f"🗑️ Removed duplicate frames from {start_ts} to {end_ts}")

    return unique_images


def create_pptx_from_images_with_timestamps(image_tuples, output_pptx, video_id):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for img_path, seconds in tqdm(image_tuples, desc="Creating slides"):
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(
            str(img_path), Inches(0), Inches(0), width=prs.slide_width
        )

        # Timestamp and hyperlink
        timestamp = f"{seconds // 60:02.0f}:{seconds % 60:02.0f}"
        youtube_link = f"https://www.youtube.com/watch?v={video_id}&t={seconds}s"

        left = Inches(0.3)
        top = prs.slide_height - Inches(0.7)  # type: ignore
        width = Inches(3)
        height = Inches(0.5)

        textbox = slide.shapes.add_textbox(left, top, width, height)  # type: ignore
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"Jump to {timestamp} on YouTube"

        font = run.font
        font.size = Pt(12)
        font.bold = True
        font.underline = True
        font.color.rgb = RGBColor(0, 102, 204)

        run.hyperlink.address = youtube_link

    prs.save(output_pptx)
    print(f"✅ PowerPoint saved: {output_pptx}")


def parse_args(argv):
    input_url_or_id = None
    custom_base = None
    fps_interval = 3  # default

    for arg in argv[1:]:
        if arg.startswith("-i=") or arg.startswith("--interval="):
            try:
                fps_interval = int(arg.split("=", 1)[1])
            except Exception:
                print("❌ Invalid interval value. Use -i=SECONDS or --interval=SECONDS")
                sys.exit(1)
        elif input_url_or_id is None:
            input_url_or_id = arg
        elif custom_base is None:
            custom_base = sanitize_filename(arg)
    return input_url_or_id, custom_base, fps_interval


# === MAIN ===
if __name__ == "__main__":
    input_url_or_id, custom_base, fps_interval = parse_args(sys.argv)
    if not input_url_or_id:
        print(
            "❌ Usage: python yt2pptx.py <YouTube_URL_or_ID> [output_base_name] [-i=SECONDS|--interval=SECONDS]"
        )
        sys.exit(1)

    out_dir = Path() / "out"
    out_dir.mkdir(exist_ok=True)

    video_file, video_title, video_id = download_youtube_video(input_url_or_id, out_dir)
    base_name = custom_base or video_title

    frames_folder = out_dir / video_id

    has_interval = False
    for p in frames_folder.glob("*_0-00-*.jpg"):
        if p.name.endswith(f"{fps_interval:02d}.jpg"):
            has_interval = True
            break

    pptx_output = out_dir / f"{base_name}.pptx"

    # Only extract frames if not already done with this interval
    if not has_interval:
        print("🎞 Extracting frames...")
        extracted_images = extract_frames_ffmpeg(
            video_file, frames_folder, interval_seconds=fps_interval
        )
    else:
        print(
            f"🎞 Using previously extracted frames in '{frames_folder}' (interval={fps_interval}s)"
        )
        extracted_images = sorted(
            Path(frames_folder).glob("frame_*.jpg"), key=sort_timestamp
        )

    print("🧹 Filtering duplicate frames...")
    unique_images = filter_unique_images(extracted_images, fps_interval=fps_interval)

    print("🧾 Creating PowerPoint...")
    create_pptx_from_images_with_timestamps(unique_images, pptx_output, video_id)
